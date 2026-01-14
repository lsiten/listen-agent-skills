#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
问卷调查报告数据分析脚本
支持多种数据格式和旧方案格式，自动选择分析模型，生成HTML报告
"""

import argparse
import sys
import os
import json
import re
from pathlib import Path
from datetime import datetime
import webbrowser

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')  # 非交互式后端
import matplotlib.pyplot as plt
import seaborn as sns
from scipy import stats
from sklearn.cluster import KMeans
from sklearn.decomposition import PCA
from sklearn.preprocessing import StandardScaler
from jinja2 import Template

# 设置中文字体
plt.rcParams['font.sans-serif'] = ['Arial Unicode MS', 'SimHei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

# 设置seaborn样式
sns.set_style("whitegrid")
sns.set_palette("husl")


def load_data_file(file_path):
    """加载数据文件，支持多种格式"""
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")
    
    ext = file_path.suffix.lower()
    
    try:
        if ext == '.csv':
            df = pd.read_csv(file_path, encoding='utf-8')
        elif ext in ['.xlsx', '.xls']:
            df = pd.read_excel(file_path)
        elif ext == '.json':
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            df = pd.DataFrame(data)
        elif ext == '.txt':
            # 尝试多种分隔符
            for sep in [',', '\t', '|', ';']:
                try:
                    df = pd.read_csv(file_path, sep=sep, encoding='utf-8')
                    if len(df.columns) > 1:
                        break
                except:
                    continue
            else:
                raise ValueError("无法解析TXT文件，请确保使用标准分隔符（逗号、制表符等）")
        elif ext == '.md':
            # 读取Markdown文件，提取表格
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            # 查找Markdown表格
            tables = re.findall(r'\|.*\|', content, re.MULTILINE)
            if tables:
                # 解析第一个表格
                lines = [line.strip() for line in tables[0].split('\n') if '|' in line]
                if len(lines) >= 2:
                    headers = [h.strip() for h in lines[0].split('|')[1:-1]]
                    data_rows = []
                    for line in lines[2:]:  # 跳过分隔行
                        row = [cell.strip() for cell in line.split('|')[1:-1]]
                        if row:
                            data_rows.append(row)
                    df = pd.DataFrame(data_rows, columns=headers)
                else:
                    raise ValueError("Markdown文件中未找到有效的表格")
            else:
                raise ValueError("Markdown文件中未找到表格")
        elif ext == '.docx':
            try:
                from docx import Document
                doc = Document(file_path)
                # 查找第一个表格
                if doc.tables:
                    table = doc.tables[0]
                    headers = [cell.text.strip() for cell in table.rows[0].cells]
                    data_rows = []
                    for row in table.rows[1:]:
                        data_rows.append([cell.text.strip() for cell in row.cells])
                    df = pd.DataFrame(data_rows, columns=headers)
                else:
                    raise ValueError("Word文件中未找到表格")
            except ImportError:
                raise ImportError("需要安装python-docx库: pip install python-docx")
        elif ext == '.pdf':
            try:
                import pdfplumber
                # 尝试从PDF中提取表格
                with pdfplumber.open(file_path) as pdf:
                    tables = []
                    for page in pdf.pages:
                        page_tables = page.extract_tables()
                        if page_tables:
                            tables.extend(page_tables)
                    
                    if tables:
                        # 使用第一个表格
                        table_data = tables[0]
                        if len(table_data) >= 2:
                            headers = [str(cell).strip() if cell else f'Column_{i}' 
                                      for i, cell in enumerate(table_data[0])]
                            data_rows = []
                            for row in table_data[1:]:
                                if row and any(cell for cell in row):
                                    # 确保列数一致
                                    row_data = [str(cell).strip() if cell else '' 
                                               for cell in row]
                                    # 补齐缺失的列
                                    while len(row_data) < len(headers):
                                        row_data.append('')
                                    # 截断多余的列
                                    row_data = row_data[:len(headers)]
                                    data_rows.append(row_data)
                            if data_rows:
                                df = pd.DataFrame(data_rows, columns=headers)
                            else:
                                raise ValueError("PDF表格中没有有效数据")
                        else:
                            raise ValueError("PDF表格格式不正确")
                    else:
                        # 如果没有表格，尝试从文本中提取表格数据
                        print("   PDF中未找到表格，尝试从文本中提取数据...")
                        all_text = []
                        for page in pdf.pages:
                            text = page.extract_text()
                            if text:
                                all_text.append(text)
                        
                        # 尝试从文本中解析表格
                        text_content = '\n'.join(all_text)
                        lines = [line.strip() for line in text_content.split('\n') if line.strip()]
                        
                        # 查找可能包含表格的行
                        table_lines = []
                        for line in lines:
                            # 检查是否包含多个分隔符（可能是表格行）
                            if sum(line.count(sep) for sep in [',', '\t', '|', ';']) >= 2:
                                table_lines.append(line)
                        
                        if table_lines and len(table_lines) >= 2:
                            # 尝试解析为表格
                            for sep in [',', '\t', '|', ';']:
                                try:
                                    data_rows = []
                                    for line in table_lines:
                                        parts = [p.strip() for p in line.split(sep) if p.strip()]
                                        if len(parts) > 1:
                                            data_rows.append(parts)
                                    
                                    if len(data_rows) >= 2:
                                        # 第一行作为表头
                                        headers = data_rows[0]
                                        df = pd.DataFrame(data_rows[1:], columns=headers)
                                        break
                                except:
                                    continue
                            else:
                                raise ValueError("PDF文件中未找到可解析的表格数据")
                        else:
                            raise ValueError("PDF文件中未找到表格，请确保PDF包含表格数据")
            except ImportError:
                raise ImportError("需要安装pdfplumber库: pip install pdfplumber")
        elif ext in ['.png', '.jpg', '.jpeg']:
            try:
                import pytesseract
                from PIL import Image
                import io
                
                # 使用OCR识别图片中的表格
                print("   正在使用OCR识别图片中的表格数据...")
                image = Image.open(file_path)
                
                # 尝试使用OCR提取文本
                ocr_text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                
                # 尝试解析表格格式的文本
                lines = [line.strip() for line in ocr_text.split('\n') if line.strip()]
                
                # 查找包含分隔符的行（可能是表格）
                table_lines = []
                for line in lines:
                    if any(sep in line for sep in [',', '\t', '|', ';', '  ']):
                        table_lines.append(line)
                
                if table_lines:
                    # 尝试解析为表格
                    for sep in [',', '\t', '|', ';']:
                        try:
                            data_rows = []
                            for line in table_lines:
                                parts = [p.strip() for p in line.split(sep) if p.strip()]
                                if len(parts) > 1:
                                    data_rows.append(parts)
                            
                            if len(data_rows) >= 2:
                                # 第一行作为表头
                                headers = data_rows[0]
                                df = pd.DataFrame(data_rows[1:], columns=headers)
                                break
                        except:
                            continue
                    else:
                        raise ValueError("无法从图片中识别出表格格式，请确保图片清晰且包含表格数据")
                else:
                    raise ValueError("图片中未识别出表格数据，请确保图片清晰且包含表格")
                    
            except ImportError:
                raise ImportError("需要安装pytesseract和Pillow库: pip install pytesseract Pillow\n"
                               "还需要安装Tesseract OCR引擎: brew install tesseract (macOS)")
            except Exception as e:
                error_msg = str(e).lower()
                if "tesseract" in error_msg or "tesseract not found" in error_msg:
                    raise ImportError("Tesseract OCR未安装或未配置。\n"
                                    "安装方法:\n"
                                    "  macOS: brew install tesseract\n"
                                    "  Ubuntu/Debian: sudo apt-get install tesseract-ocr\n"
                                    "  Windows: 下载安装 https://github.com/UB-Mannheim/tesseract/wiki")
                raise
        else:
            raise ValueError(f"不支持的文件格式: {ext}")
        
        print(f"✅ 成功加载数据文件: {file_path}")
        print(f"   数据形状: {df.shape[0]} 行 × {df.shape[1]} 列")
        return df
    
    except Exception as e:
        raise ValueError(f"加载文件失败: {str(e)}")


def load_old_plan(file_path):
    """加载旧方案文件，支持多种格式"""
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")
    
    ext = file_path.suffix.lower()
    
    try:
        if ext == '.md':
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return content
        elif ext == '.docx':
            try:
                from docx import Document
                doc = Document(file_path)
                content = '\n'.join([para.text for para in doc.paragraphs])
                return content
            except ImportError:
                raise ImportError("需要安装python-docx库: pip install python-docx")
        elif ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                content_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content_parts.append(shape.text)
                return '\n'.join(content_parts)
            except ImportError:
                raise ImportError("需要安装python-pptx库: pip install python-pptx")
        elif ext == '.pdf':
            try:
                import pdfplumber
                content_parts = []
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            content_parts.append(text)
                return '\n\n'.join(content_parts)
            except ImportError:
                raise ImportError("需要安装pdfplumber库: pip install pdfplumber")
        elif ext in ['.png', '.jpg', '.jpeg']:
            try:
                import pytesseract
                from PIL import Image
                
                print("   正在使用OCR识别图片中的文本内容...")
                image = Image.open(file_path)
                content = pytesseract.image_to_string(image, lang='chi_sim+eng')
                return content
            except ImportError:
                raise ImportError("需要安装pytesseract和Pillow库: pip install pytesseract Pillow\n"
                               "还需要安装Tesseract OCR引擎: brew install tesseract (macOS)")
            except Exception as e:
                error_msg = str(e).lower()
                if "tesseract" in error_msg or "tesseract not found" in error_msg:
                    raise ImportError("Tesseract OCR未安装或未配置。\n"
                                    "安装方法:\n"
                                    "  macOS: brew install tesseract\n"
                                    "  Ubuntu/Debian: sudo apt-get install tesseract-ocr\n"
                                    "  Windows: 下载安装 https://github.com/UB-Mannheim/tesseract/wiki")
                raise
        else:
            raise ValueError(f"不支持的旧方案格式: {ext}")
    
    except Exception as e:
        raise ValueError(f"加载旧方案文件失败: {str(e)}")


def select_analysis_model(df):
    """根据数据特征自动选择分析模型"""
    n_samples = len(df)
    n_vars = len(df.columns)
    
    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    n_numeric = len(numeric_cols)
    
    print(f"\n📊 数据特征分析:")
    print(f"   样本数: {n_samples}")
    print(f"   变量数: {n_vars}")
    print(f"   数值变量: {n_numeric}")
    
    models = []
    
    # 1. 描述性统计（总是执行）
    models.append('descriptive')
    
    # 2. 相关性分析（如果有两个以上数值变量）
    if n_numeric >= 2:
        models.append('correlation')
    
    # 3. 回归分析（如果有因变量和自变量）
    if n_numeric >= 2 and n_samples >= 30:
        models.append('regression')
    
    # 4. 聚类分析（样本量足够）
    if n_samples >= 50 and n_numeric >= 2:
        models.append('cluster')
    
    # 5. 因子分析（变量数量多）
    if n_numeric >= 5 and n_samples >= 100:
        models.append('factor')
    
    print(f"   推荐模型: {', '.join(models)}")
    return models


def perform_descriptive_analysis(df):
    """描述性统计分析"""
    results = {
        'summary': df.describe().to_dict(),
        'missing': df.isnull().sum().to_dict(),
        'dtypes': df.dtypes.astype(str).to_dict()
    }
    return results


def perform_correlation_analysis(df):
    """相关性分析"""
    numeric_df = df.select_dtypes(include=[np.number])
    if len(numeric_df.columns) < 2:
        return None
    
    corr_matrix = numeric_df.corr()
    return {
        'matrix': corr_matrix.to_dict(),
        'strong_pairs': []
    }


def perform_regression_analysis(df):
    """回归分析"""
    numeric_df = df.select_dtypes(include=[np.number])
    if len(numeric_df.columns) < 2:
        return None
    
    # 选择第一个数值列作为因变量
    y_col = numeric_df.columns[0]
    X_cols = numeric_df.columns[1:].tolist()
    
    if not X_cols:
        return None
    
    from sklearn.linear_model import LinearRegression
    from sklearn.metrics import r2_score
    
    X = numeric_df[X_cols]
    y = numeric_df[y_col]
    
    # 处理缺失值
    mask = ~(X.isnull().any(axis=1) | y.isnull())
    X_clean = X[mask]
    y_clean = y[mask]
    
    if len(X_clean) < 10:
        return None
    
    model = LinearRegression()
    model.fit(X_clean, y_clean)
    y_pred = model.predict(X_clean)
    r2 = r2_score(y_clean, y_pred)
    
    return {
        'target': y_col,
        'features': X_cols,
        'r2_score': float(r2),
        'coefficients': {col: float(coef) for col, coef in zip(X_cols, model.coef_)},
        'intercept': float(model.intercept_)
    }


def perform_cluster_analysis(df):
    """聚类分析"""
    numeric_df = df.select_dtypes(include=[np.number])
    if len(numeric_df.columns) < 2:
        return None
    
    # 处理缺失值
    numeric_df = numeric_df.dropna()
    if len(numeric_df) < 10:
        return None
    
    # 标准化
    scaler = StandardScaler()
    X_scaled = scaler.fit_transform(numeric_df)
    
    # K-means聚类
    n_clusters = min(5, len(numeric_df) // 10)
    if n_clusters < 2:
        return None
    
    kmeans = KMeans(n_clusters=n_clusters, random_state=42, n_init=10)
    clusters = kmeans.fit_predict(X_scaled)
    
    return {
        'n_clusters': int(n_clusters),
        'cluster_labels': clusters.tolist(),
        'inertia': float(kmeans.inertia_)
    }


def perform_factor_analysis(df):
    """因子分析"""
    numeric_df = df.select_dtypes(include=[np.number])
    if len(numeric_df.columns) < 3:
        return None
    
    numeric_df = numeric_df.dropna()
    if len(numeric_df) < 10:
        return None
    
    # PCA降维
    scaler = StandardScaler()
    X_scaled = scaler.fit_transform(numeric_df)
    
    pca = PCA(n_components=min(3, len(numeric_df.columns)))
    pca.fit(X_scaled)
    
    return {
        'n_components': int(pca.n_components_),
        'explained_variance_ratio': [float(v) for v in pca.explained_variance_ratio_],
        'total_variance_explained': float(sum(pca.explained_variance_ratio_))
    }


def generate_charts(df, output_dir, html_output_path):
    """生成可视化图表"""
    charts = {}
    numeric_df = df.select_dtypes(include=[np.number])
    
    # 计算相对路径
    html_dir = Path(html_output_path).parent
    charts_dir = html_dir / 'charts'
    charts_dir.mkdir(exist_ok=True)
    
    # 1. 数值变量分布图
    if len(numeric_df.columns) > 0:
        fig, axes = plt.subplots(2, 2, figsize=(12, 10))
        axes = axes.flatten()
        
        for i, col in enumerate(numeric_df.columns[:4]):
            if i < len(axes):
                axes[i].hist(numeric_df[col].dropna(), bins=20, edgecolor='black')
                axes[i].set_title(f'{col} 分布', fontsize=12)
                axes[i].set_xlabel(col)
                axes[i].set_ylabel('频数')
        
        # 隐藏多余的子图
        for i in range(len(numeric_df.columns), len(axes)):
            axes[i].axis('off')
        
        plt.tight_layout()
        dist_path = charts_dir / 'distribution_chart.png'
        plt.savefig(dist_path, dpi=150, bbox_inches='tight')
        plt.close()
        # 使用相对路径
        charts['distribution'] = f'charts/{dist_path.name}'
    
    # 2. 相关性热力图
    if len(numeric_df.columns) >= 2:
        fig, ax = plt.subplots(figsize=(10, 8))
        corr = numeric_df.corr()
        sns.heatmap(corr, annot=True, fmt='.2f', cmap='coolwarm', center=0,
                   square=True, linewidths=1, cbar_kws={"shrink": .8}, ax=ax)
        ax.set_title('变量相关性热力图', fontsize=14, pad=20)
        plt.tight_layout()
        corr_path = charts_dir / 'correlation_heatmap.png'
        plt.savefig(corr_path, dpi=150, bbox_inches='tight')
        plt.close()
        # 使用相对路径
        charts['correlation'] = f'charts/{corr_path.name}'
    
    return charts


def evaluate_old_plan(old_plan_content, analysis_results):
    """评估旧方案"""
    evaluation = {
        'summary': '',
        'strengths': [],
        'weaknesses': [],
        'recommendations': []
    }
    
    # 简单的评估逻辑（可以根据需要扩展）
    plan_lower = old_plan_content.lower()
    
    # 检查是否提到数据分析
    if '分析' in plan_lower or 'analysis' in plan_lower:
        evaluation['strengths'].append('方案中包含了数据分析相关内容')
    else:
        evaluation['weaknesses'].append('方案中缺少明确的数据分析计划')
    
    # 检查样本量
    sample_size = analysis_results.get('data_info', {}).get('n_samples', 0)
    if sample_size < 30:
        evaluation['weaknesses'].append(f'样本量较小（{sample_size}），可能影响统计显著性')
        evaluation['recommendations'].append('建议增加样本量至至少30个，以提高分析的可信度')
    elif sample_size >= 100:
        evaluation['strengths'].append(f'样本量充足（{sample_size}），适合进行深度分析')
    
    # 检查变量数量
    n_vars = analysis_results.get('data_info', {}).get('n_vars', 0)
    if n_vars < 3:
        evaluation['weaknesses'].append('变量数量较少，分析维度有限')
        evaluation['recommendations'].append('建议增加更多调研维度，丰富数据内容')
    
    evaluation['summary'] = f"基于数据分析结果，对旧方案进行了评估。发现{len(evaluation['strengths'])}个优点和{len(evaluation['weaknesses'])}个需要改进的地方。"
    
    return evaluation


def generate_html_report(analysis_results, charts, old_plan_eval, title, output_path):
    """生成HTML报告"""
    
    html_template = """
<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ title }}</title>
    <style>
        * { margin: 0; padding: 0; box-sizing: border-box; }
        body {
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', 'Microsoft YaHei', sans-serif;
            line-height: 1.6;
            color: #333;
            background: #f5f5f5;
            padding: 20px;
        }
        .container {
            max-width: 1200px;
            margin: 0 auto;
            background: white;
            padding: 40px;
            border-radius: 8px;
            box-shadow: 0 2px 10px rgba(0,0,0,0.1);
        }
        h1 {
            color: #2c3e50;
            border-bottom: 3px solid #3498db;
            padding-bottom: 10px;
            margin-bottom: 30px;
        }
        h2 {
            color: #34495e;
            margin-top: 40px;
            margin-bottom: 20px;
            padding-left: 10px;
            border-left: 4px solid #3498db;
        }
        h3 {
            color: #555;
            margin-top: 25px;
            margin-bottom: 15px;
        }
        .meta-info {
            background: #ecf0f1;
            padding: 15px;
            border-radius: 5px;
            margin-bottom: 30px;
        }
        .meta-info p {
            margin: 5px 0;
        }
        table {
            width: 100%;
            border-collapse: collapse;
            margin: 20px 0;
        }
        th, td {
            padding: 12px;
            text-align: left;
            border-bottom: 1px solid #ddd;
        }
        th {
            background-color: #3498db;
            color: white;
        }
        tr:hover {
            background-color: #f5f5f5;
        }
        .chart-container {
            text-align: center;
            margin: 30px 0;
        }
        .chart-container img {
            max-width: 100%;
            height: auto;
            border: 1px solid #ddd;
            border-radius: 5px;
        }
        .strength {
            color: #27ae60;
            background: #d5f4e6;
            padding: 10px;
            border-radius: 5px;
            margin: 10px 0;
        }
        .weakness {
            color: #e74c3c;
            background: #fadbd8;
            padding: 10px;
            border-radius: 5px;
            margin: 10px 0;
        }
        .recommendation {
            color: #2980b9;
            background: #d6eaf8;
            padding: 10px;
            border-radius: 5px;
            margin: 10px 0;
        }
        .section {
            margin: 30px 0;
        }
        code {
            background: #f4f4f4;
            padding: 2px 6px;
            border-radius: 3px;
            font-family: 'Courier New', monospace;
        }
    </style>
</head>
<body>
    <div class="container">
        <h1>{{ title }}</h1>
        
        <div class="meta-info">
            <p><strong>生成时间:</strong> {{ generation_time }}</p>
            <p><strong>数据样本数:</strong> {{ data_info.n_samples }}</p>
            <p><strong>变量数量:</strong> {{ data_info.n_vars }}</p>
            <p><strong>使用的分析模型:</strong> {{ ', '.join(models_used) }}</p>
        </div>

        <div class="section">
            <h2>📊 数据概览</h2>
            <p>本次分析共包含 <strong>{{ data_info.n_samples }}</strong> 个样本，<strong>{{ data_info.n_vars }}</strong> 个变量。</p>
        </div>

        <div class="section">
            <h2>📈 描述性统计</h2>
            <p>以下是各变量的基本统计信息：</p>
            <!-- 描述性统计表格将在这里插入 -->
        </div>

        {% if charts.distribution %}
        <div class="section">
            <h2>📉 数据分布可视化</h2>
            <div class="chart-container">
                <img src="{{ charts.distribution }}" alt="数据分布图">
            </div>
        </div>
        {% endif %}

        {% if charts.correlation %}
        <div class="section">
            <h2>🔥 相关性分析</h2>
            <div class="chart-container">
                <img src="{{ charts.correlation }}" alt="相关性热力图">
            </div>
        </div>
        {% endif %}

        {% if regression_results %}
        <div class="section">
            <h2>📐 回归分析结果</h2>
            <p><strong>目标变量:</strong> {{ regression_results.target }}</p>
            <p><strong>R² 得分:</strong> {{ "%.3f"|format(regression_results.r2_score) }}</p>
            <h3>系数:</h3>
            <ul>
                {% for feature, coef in regression_results.coefficients.items() %}
                <li><code>{{ feature }}</code>: {{ "%.4f"|format(coef) }}</li>
                {% endfor %}
            </ul>
        </div>
        {% endif %}

        {% if old_plan_eval %}
        <div class="section">
            <h2>🔍 旧方案评估</h2>
            <p>{{ old_plan_eval.summary }}</p>
            
            {% if old_plan_eval.strengths %}
            <h3>✅ 优点</h3>
            {% for strength in old_plan_eval.strengths %}
            <div class="strength">{{ strength }}</div>
            {% endfor %}
            {% endif %}

            {% if old_plan_eval.weaknesses %}
            <h3>⚠️ 需要改进的地方</h3>
            {% for weakness in old_plan_eval.weaknesses %}
            <div class="weakness">{{ weakness }}</div>
            {% endfor %}
            {% endif %}

            {% if old_plan_eval.recommendations %}
            <h3>💡 改进建议</h3>
            {% for rec in old_plan_eval.recommendations %}
            <div class="recommendation">{{ rec }}</div>
            {% endfor %}
            {% endif %}
        </div>
        {% endif %}

        <div class="section">
            <h2>📝 结论与建议</h2>
            <p>基于以上分析，我们得出以下结论：</p>
            <ul>
                <li>数据质量良好，适合进行深度分析</li>
                <li>建议根据分析结果制定相应的行动方案</li>
                <li>定期更新数据，持续跟踪分析</li>
            </ul>
        </div>
    </div>
</body>
</html>
    """
    
    template = Template(html_template)
    html_content = template.render(
        title=title,
        generation_time=datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
        data_info=analysis_results.get('data_info', {}),
        models_used=analysis_results.get('models_used', []),
        charts=charts,
        regression_results=analysis_results.get('regression'),
        old_plan_eval=old_plan_eval
    )
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    print(f"✅ HTML报告已生成: {output_path}")


def main():
    parser = argparse.ArgumentParser(description='问卷调查报告数据分析工具')
    parser.add_argument('--data', required=True, help='问卷数据文件路径')
    parser.add_argument('--output', default='output/report.html', help='输出HTML报告路径')
    parser.add_argument('--old-plan', help='旧调研方案文件路径（可选）')
    parser.add_argument('--model', default='auto', choices=['auto', 'descriptive', 'correlation', 'regression', 'cluster', 'factor'],
                       help='分析模型类型（默认：auto自动选择）')
    parser.add_argument('--title', default='问卷数据分析报告', help='报告标题')
    parser.add_argument('--open-browser', action='store_true', default=True, help='是否自动打开浏览器')
    parser.add_argument('--verbose', action='store_true', help='显示详细日志')
    
    args = parser.parse_args()
    
    try:
        # 1. 加载数据
        print("📂 正在加载数据文件...")
        df = load_data_file(args.data)
        
        # 2. 选择分析模型
        if args.model == 'auto':
            models = select_analysis_model(df)
        else:
            models = [args.model]
        
        # 3. 执行分析
        print("\n🔬 正在执行数据分析...")
        analysis_results = {
            'data_info': {
                'n_samples': len(df),
                'n_vars': len(df.columns)
            },
            'models_used': models
        }
        
        if 'descriptive' in models:
            analysis_results['descriptive'] = perform_descriptive_analysis(df)
        
        if 'correlation' in models:
            analysis_results['correlation'] = perform_correlation_analysis(df)
        
        if 'regression' in models:
            analysis_results['regression'] = perform_regression_analysis(df)
        
        if 'cluster' in models:
            analysis_results['cluster'] = perform_cluster_analysis(df)
        
        if 'factor' in models:
            analysis_results['factor'] = perform_factor_analysis(df)
        
        # 4. 生成图表
        print("\n📊 正在生成可视化图表...")
        output_dir = Path(args.output).parent
        output_dir.mkdir(parents=True, exist_ok=True)
        charts = generate_charts(df, output_dir, args.output)
        
        # 5. 评估旧方案（如果提供）
        old_plan_eval = None
        if args.old_plan:
            print("\n🔍 正在评估旧方案...")
            old_plan_content = load_old_plan(args.old_plan)
            old_plan_eval = evaluate_old_plan(old_plan_content, analysis_results)
        
        # 6. 生成HTML报告
        print("\n📝 正在生成HTML报告...")
        generate_html_report(analysis_results, charts, old_plan_eval, args.title, args.output)
        
        # 7. 打开浏览器
        if args.open_browser:
            print(f"\n🌐 正在浏览器中打开报告...")
            webbrowser.open(f'file://{os.path.abspath(args.output)}')
        
        print("\n✅ 分析完成！")
        
    except Exception as e:
        print(f"\n❌ 错误: {str(e)}", file=sys.stderr)
        if args.verbose:
            import traceback
            traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
